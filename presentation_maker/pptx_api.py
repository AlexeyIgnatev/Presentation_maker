from pptx import Presentation
from presentation_maker.saver import check_path, path_to_file
from presentation_maker.config import source_path
from pptx.util import Inches
import os


class Pptx:
    def __init__(self):
        self.prs = None
        self.topic = ''
        self.file = ''

    def create_presentation(self, topic, headers, image_paths):
        self.topic = topic
        self.prs = Presentation()
        title_slide_layout = self.prs.slide_layouts[0]
        slide = self.prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        title.text = self.topic

        title_only_slide_layout = self.prs.slide_layouts[5]

        for head, images in zip(headers, image_paths[0].values()):
            slide = self.prs.slides.add_slide(title_only_slide_layout)
            shapes = slide.shapes
            title_shape = shapes.title

            left = top = Inches(1.5)
            width = Inches(7)
            height = Inches(5)
            if images[0]:
                pic = shapes.add_picture(images[0], left, top, width=width, height=height)
            title_shape.text = head

    def save_presentation(self):
        if self.prs is None:
            print('Создайте презентацию перед соханением!')
        else:
            check_path(source_path)
            check_path('{}/{}'.format(source_path, self.topic))
            self.file = '{}/{}/{}.pptx'.format(source_path, self.topic, self.topic)
            self.prs.save(self.file)

    def open_presentation(self):
        if self.file:
            os.startfile(path_to_file(self.file))
